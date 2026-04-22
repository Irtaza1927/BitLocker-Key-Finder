"""
partition_scan.py
─────────────────────────────────────────────────────────────────────────────
BitLocker Key Finder — Core Scanning Engine (Part B)
FAST-NUCES Islamabad | Digital Forensics | Semester 6

Fixes in this revision:
  [F1] Size filter (_ACTIVE_MIN_BYTES / _ACTIVE_MAX_BYTES) now actually
       applied in scan_partition() before parser dispatch.
  [F2] Extension filter (_ACTIVE_EXTS) now actually applied in
       scan_partition().
  [F3] Report timestamp computed inside save_report() (no longer frozen
       at module-import time).
  [F4] Dead code removed from parse_bek (meaningless binary GUID regex).
  [F5] Multiple keys in a single file are now paired with their
       spatially-closest Recovery-Key-ID via regex match offsets
       (previously every key got ids[0]).
  [F6] BitLocker mod-11 structural validation added. Every 6-digit
       block of a 48-digit candidate must satisfy:
             block % 11 == 0   AND   block // 11 <= 0xFFFF
       This is Microsoft's on-disk BitLocker recovery-password invariant
       (ref: dislocker / libbde documentation).
       Pattern-only matches are still reported but marked as such, so
       forensic evidence is not discarded.

Supported file types : txt, docx, xlsx, pptx, pdf, rtf, eml, odt, bek,
                       plus images / db / vhd / iso (raw carve).
"""

from __future__ import annotations

import re
import os
import csv
import datetime
import chardet
import fnmatch

# ─── Regex Patterns ──────────────────────────────────────────────────────────
# 48-digit BitLocker recovery password: 8 groups of 6 digits
KEY_PATTERN = re.compile(
    r"\b\d{6}-\d{6}-\d{6}-\d{6}-\d{6}-\d{6}-\d{6}-\d{6}\b"
)
# Recovery Key Identification GUID
ID_PATTERN = re.compile(
    r"\b[0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{12}\b"
)

# ─── Runtime config (written by GUI, read by scanner) ────────────────────────
_ACTIVE_MIN_BYTES: int = 0
_ACTIVE_MAX_BYTES: int = 50 * 1024 * 1024   # 50 MB default
_ACTIVE_EXTS: set | None = None             # None → accept all mapped exts

# Hard absolute ceiling — protects against accidental 10 GB reads even
# if the user sets Max = "999 GB". Raise at your own risk.
_ABSOLUTE_SIZE_CEILING = 500 * 1024 * 1024  # 500 MB

# ─── File Extension Categories ───────────────────────────────────────────────
TEXT_EXTENSIONS = {
    ".txt", ".log", ".md", ".ini", ".cfg", ".conf", ".bat", ".ps1",
    ".csv", ".tsv", ".json", ".xml", ".html", ".htm", ".yaml", ".yml",
    ".bek", ".key", ".asc", ".pem", ".crt", ".cer", ".sql", ".py",
    ".js", ".ts", ".css", ".sh", ".reg",
}
BINARY_TEXT_EXTENSIONS = {".docx", ".xlsx", ".pptx", ".odt", ".ods", ".odp"}
SPECIAL_EXTENSIONS = {".pdf", ".rtf", ".eml", ".msg", ".mdb", ".accdb"}


# ═════════════════════════════════════════════════════════════════════════════
#  VALIDATION
# ═════════════════════════════════════════════════════════════════════════════
def validate_bitlocker_key(key_str: str) -> bool:
    """
    Structural validation of a 48-digit BitLocker recovery password.

    Microsoft's format stores the 128-bit stretch key as 8 * 16-bit integers.
    Each 16-bit integer `q` is printed as a 6-digit decimal block = q * 11.
    Therefore a genuine block must satisfy both:
        block % 11 == 0
        block // 11 <= 0xFFFF

    This mod-11 check rejects ~90% of false-positive pattern matches
    (random 48-digit strings in logs, serial numbers, etc.).
    """
    groups = key_str.split("-")
    if len(groups) != 8:
        return False
    for g in groups:
        if len(g) != 6 or not g.isdigit():
            return False
        n = int(g)
        if n % 11 != 0:
            return False
        if n // 11 > 0xFFFF:
            return False
    return True


# ═════════════════════════════════════════════════════════════════════════════
#  INTERNAL HELPERS
# ═════════════════════════════════════════════════════════════════════════════
def _size_allowed(filepath: str) -> bool:
    """Respect GUI-configured min/max range plus an absolute safety ceiling."""
    try:
        size = os.path.getsize(filepath)
    except OSError:
        return False
    if size < _ACTIVE_MIN_BYTES:
        return False
    effective_max = min(_ACTIVE_MAX_BYTES, _ABSOLUTE_SIZE_CEILING)
    if size > effective_max:
        return False
    return True


def _pair_keys_with_ids(text: str):
    """
    Yield (key_str, id_str_or_None) pairs using closest-offset matching.

    For each key match, the Recovery-Key-ID chosen is the one whose
    regex-match start offset is nearest to the key's offset.
    Falls back to None if no IDs exist in the text.
    """
    key_matches = list(KEY_PATTERN.finditer(text))
    id_matches  = list(ID_PATTERN.finditer(text))
    for km in key_matches:
        if id_matches:
            nearest = min(id_matches, key=lambda im: abs(im.start() - km.start()))
            yield km.group(), nearest.group()
        else:
            yield km.group(), None


def _extract_from_text(text: str, filepath: str, results: list, log_fn=None) -> bool:
    """
    Core extraction: scans decoded text for BitLocker keys + IDs,
    applies mod-11 validation, appends structured records to `results`.
    Returns True if at least one key (valid or pattern-only) was found.
    """
    found_any = False
    for key_str, rid in _pair_keys_with_ids(text):
        is_valid = validate_bitlocker_key(key_str)
        record = {
            "File Path":       filepath,
            "Recovery Key ID": rid or "Unknown",
            "BitLocker Key":   key_str,
            "Validation":      "Valid (mod-11)" if is_valid else "Pattern-only",
            "Source":          "Text Content",
        }
        results.append(record)
        found_any = True
        if log_fn:
            tag = "success" if is_valid else "warning"
            label = "VALID KEY" if is_valid else "PATTERN MATCH"
            log_fn(f"[{label}] {key_str}  |  ID: {rid or 'Unknown'}  "
                   f"|  File: {filepath}", tag)
    return found_any


def _decode_best_effort(raw: bytes):
    """Yield successive decode attempts: BOM-aware first, then chardet."""
    for enc in ("utf-8-sig", "utf-8", "utf-16", "utf-16-le",
                "utf-16-be", "latin-1"):
        try:
            yield raw.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    # chardet fallback
    try:
        detected = chardet.detect(raw) or {}
        enc = detected.get("encoding")
        if enc:
            yield raw.decode(enc, errors="replace")
    except Exception:
        pass


# ═════════════════════════════════════════════════════════════════════════════
#  INDIVIDUAL PARSERS
# ═════════════════════════════════════════════════════════════════════════════
def parse_txt(filepath, results, log_fn=None):
    """Plain-text family: read raw bytes, try multiple encodings."""
    try:
        with open(filepath, "rb") as f:
            raw = f.read()
        for text in _decode_best_effort(raw):
            if _extract_from_text(text, filepath, results, log_fn):
                return
    except PermissionError:
        if log_fn:
            log_fn(f"[PERMISSION DENIED] {filepath}", "warning")
    except Exception as e:
        if log_fn:
            log_fn(f"[ERROR] {filepath}: {e}", "error")


def parse_docx(filepath, results, log_fn=None):
    try:
        import docx
        doc = docx.Document(filepath)
        combined = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    combined.append(cell.text)
        try:
            props = doc.core_properties
            for attr in ("comments", "subject", "title", "description"):
                combined.append(getattr(props, attr, "") or "")
        except Exception:
            pass
        _extract_from_text("\n".join(combined), filepath, results, log_fn)
    except ImportError:
        if log_fn:
            log_fn("[MISSING LIB] python-docx not installed. "
                   "Run: pip install python-docx", "error")
    except Exception as e:
        if log_fn:
            log_fn(f"[ERROR] {filepath}: {e}", "error")


def parse_xlsx(filepath, results, log_fn=None):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        combined = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                combined.append(" | ".join(str(c) for c in row if c is not None))
        wb.close()
        _extract_from_text("\n".join(combined), filepath, results, log_fn)
    except ImportError:
        try:
            import pandas as pd
            xls = pd.ExcelFile(filepath)
            combined = [xls.parse(s).to_string() for s in xls.sheet_names]
            _extract_from_text("\n".join(combined), filepath, results, log_fn)
        except Exception as e2:
            if log_fn:
                log_fn(f"[ERROR] xlsx {filepath}: {e2}", "error")
    except Exception as e:
        if log_fn:
            log_fn(f"[ERROR] {filepath}: {e}", "error")


def parse_pptx(filepath, results, log_fn=None):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        combined = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    combined.append(shape.text or "")
                if getattr(shape, "has_text_frame", False):
                    for para in shape.text_frame.paragraphs:
                        combined.append(para.text or "")
        _extract_from_text("\n".join(combined), filepath, results, log_fn)
    except ImportError:
        if log_fn:
            log_fn("[MISSING LIB] python-pptx not installed. "
                   "Run: pip install python-pptx", "warning")
    except Exception as e:
        if log_fn:
            log_fn(f"[ERROR] {filepath}: {e}", "error")


def parse_pdf(filepath, results, log_fn=None):
    # Primary: pdfplumber
    try:
        import pdfplumber
        combined = []
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                combined.append(page.extract_text() or "")
                for table in page.extract_tables() or []:
                    for row in table:
                        combined.append(" ".join(str(c) for c in row if c))
        _extract_from_text("\n".join(combined), filepath, results, log_fn)
        return
    except ImportError:
        pass
    except Exception as e:
        if log_fn:
            log_fn(f"[WARN] pdfplumber failed: {e}", "warning")

    # Fallback: PyPDF2
    try:
        import PyPDF2
        combined = []
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                combined.append(page.extract_text() or "")
        _extract_from_text("\n".join(combined), filepath, results, log_fn)
        return
    except ImportError:
        pass
    except Exception as e:
        if log_fn:
            log_fn(f"[WARN] PyPDF2 failed: {e}", "warning")

    # Last resort: raw carve
    parse_raw_binary(filepath, results, log_fn)


def parse_rtf(filepath, results, log_fn=None):
    try:
        from striprtf.striprtf import rtf_to_text
        with open(filepath, "rb") as f:
            raw = f.read()
        for enc in ("utf-8", "latin-1"):
            try:
                text = rtf_to_text(raw.decode(enc, errors="replace"))
                _extract_from_text(text, filepath, results, log_fn)
                return
            except Exception:
                continue
    except ImportError:
        try:
            with open(filepath, "rb") as f:
                raw = f.read()
            text = re.sub(rb"\\[a-z]+\d*\s?|[{}]", b" ", raw)\
                     .decode("latin-1", errors="replace")
            _extract_from_text(text, filepath, results, log_fn)
        except Exception as e:
            if log_fn:
                log_fn(f"[ERROR] RTF {filepath}: {e}", "error")
    except Exception as e:
        if log_fn:
            log_fn(f"[ERROR] {filepath}: {e}", "error")


def parse_eml(filepath, results, log_fn=None):
    try:
        import email
        from email import policy
        with open(filepath, "rb") as f:
            msg = email.message_from_binary_file(f, policy=policy.default)
        combined = [str(msg.get("Subject", ""))]
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() in ("text/plain", "text/html"):
                    try:
                        combined.append(part.get_content())
                    except Exception:
                        pass
        else:
            try:
                combined.append(msg.get_content())
            except Exception:
                pass
        _extract_from_text("\n".join(combined), filepath, results, log_fn)
    except Exception as e:
        if log_fn:
            log_fn(f"[ERROR] EML {filepath}: {e}", "error")


def parse_odt(filepath, results, log_fn=None):
    try:
        import zipfile
        combined = []
        with zipfile.ZipFile(filepath, "r") as z:
            if "content.xml" in z.namelist():
                xml_content = z.read("content.xml").decode("utf-8", errors="replace")
                combined.append(re.sub(r"<[^>]+>", " ", xml_content))
        _extract_from_text("\n".join(combined), filepath, results, log_fn)
    except Exception as e:
        if log_fn:
            log_fn(f"[ERROR] ODT {filepath}: {e}", "error")


def parse_raw_binary(filepath, results, log_fn=None):
    """Carve raw binary: scan bytes + multi-encoding decodes."""
    try:
        with open(filepath, "rb") as f:
            raw = f.read()
        # Decoded-text pass (covers UTF-16-encoded strings in memory dumps)
        for text in _decode_best_effort(raw):
            _extract_from_text(text, filepath, results, log_fn)
        # Direct byte-pattern pass (catches ASCII in binaries that fail all decodes)
        ascii_pat = re.compile(
            rb"\b\d{6}-\d{6}-\d{6}-\d{6}-\d{6}-\d{6}-\d{6}-\d{6}\b"
        )
        for m in ascii_pat.finditer(raw):
            key_str = m.group().decode("ascii")
            # Dedup against earlier passes
            if any(r["BitLocker Key"] == key_str and r["File Path"] == filepath
                   for r in results):
                continue
            is_valid = validate_bitlocker_key(key_str)
            results.append({
                "File Path":       filepath,
                "Recovery Key ID": "Unknown (Carved)",
                "BitLocker Key":   key_str,
                "Validation":      "Valid (mod-11)" if is_valid else "Pattern-only",
                "Source":          "Raw Binary Carving",
            })
            if log_fn:
                tag = "success" if is_valid else "warning"
                log_fn(f"[CARVED {('VALID' if is_valid else 'PATTERN')}] {key_str}  "
                       f"|  File: {filepath}", tag)
    except PermissionError:
        if log_fn:
            log_fn(f"[PERMISSION DENIED] {filepath}", "warning")
    except Exception as e:
        if log_fn:
            log_fn(f"[ERROR] Binary scan {filepath}: {e}", "error")


# ─── BEK File Handler ────────────────────────────────────────────────────────
def parse_bek(filepath, results, log_fn=None):
    """
    BitLocker External Key (.bek). Binary format; carries a GUID identifier
    in UTF-16-LE near the header. The raw VMK ciphertext cannot be
    represented as a recovery password (different key type).
    """
    try:
        with open(filepath, "rb") as f:
            raw = f.read()
        rid = None
        for text in _decode_best_effort(raw):
            m = ID_PATTERN.search(text)
            if m:
                rid = m.group()
                break
        results.append({
            "File Path":       filepath,
            "Recovery Key ID": rid or "Unknown",
            "BitLocker Key":   "Binary Key Material (BEK — VMK ciphertext)",
            "Validation":      "N/A (binary container)",
            "Source":          "BEK File",
        })
        if log_fn:
            log_fn(f"[BEK FILE] {filepath}  |  ID: {rid or 'Unknown'}", "success")
    except PermissionError:
        if log_fn:
            log_fn(f"[PERMISSION DENIED] {filepath}", "warning")
    except Exception as e:
        if log_fn:
            log_fn(f"[ERROR] BEK {filepath}: {e}", "error")


# ═════════════════════════════════════════════════════════════════════════════
#  EXTENSION → PARSER MAP
# ═════════════════════════════════════════════════════════════════════════════
EXTENSION_PARSER_MAP = {
    # Text-based
    ".txt":   parse_txt, ".log":  parse_txt, ".md":   parse_txt,
    ".ini":   parse_txt, ".cfg":  parse_txt, ".conf": parse_txt,
    ".bat":   parse_txt, ".ps1":  parse_txt, ".tsv":  parse_txt,
    ".json":  parse_txt, ".xml":  parse_txt, ".html": parse_txt,
    ".htm":   parse_txt, ".yaml": parse_txt, ".yml":  parse_txt,
    ".key":   parse_txt, ".asc":  parse_txt, ".pem":  parse_txt,
    ".crt":   parse_txt, ".cer":  parse_txt, ".sql":  parse_txt,
    ".py":    parse_txt, ".js":   parse_txt, ".ts":   parse_txt,
    ".css":   parse_txt, ".sh":   parse_txt, ".reg":  parse_txt,
    ".csv":   parse_txt,
    # Office
    ".docx":  parse_docx,
    ".xlsx":  parse_xlsx,
    ".pptx":  parse_pptx,
    ".odt":   parse_odt, ".ods": parse_odt, ".odp": parse_odt,
    # PDF / RTF / email
    ".pdf":   parse_pdf,
    ".rtf":   parse_rtf,
    ".eml":   parse_eml,
    # BitLocker External Key
    ".bek":   parse_bek,
    # Raw carve targets
    ".png":   parse_raw_binary, ".jpg":   parse_raw_binary,
    ".jpeg":  parse_raw_binary, ".bmp":   parse_raw_binary,
    ".gif":   parse_raw_binary, ".db":    parse_raw_binary,
    ".sqlite": parse_raw_binary, ".mdb":  parse_raw_binary,
    ".accdb": parse_raw_binary, ".dat":   parse_raw_binary,
    ".bin":   parse_raw_binary, ".img":   parse_raw_binary,
    ".iso":   parse_raw_binary, ".vhd":   parse_raw_binary,
    ".vhdx":  parse_raw_binary, ".vmdk":  parse_raw_binary,
}

BITLOCKER_NAME_PATTERNS = [
    "*BitLocker Recovery Key*",
    "*bitlocker*",
    "*recovery key*",
    "*RecoveryKey*",
    "*.bek",
    "*.BEK",
]


# ═════════════════════════════════════════════════════════════════════════════
#  MAIN SCANNER
# ═════════════════════════════════════════════════════════════════════════════
def scan_partition(folder, results, log_fn=None, progress_fn=None,
                   do_name_search=True, do_content_search=True,
                   do_binary_carve=False, stop_flag=None):
    """
    Walk `folder` recursively; dispatch each file to its parser.

    Respects module-level runtime config:
        _ACTIVE_MIN_BYTES / _ACTIVE_MAX_BYTES  (file-size filter)
        _ACTIVE_EXTS                           (extension filter, or None)
    """
    all_files = []
    if log_fn:
        log_fn(f"[SCAN] Walking directory: {folder}", "info")

    # Warn once if not elevated — user needs to know why some dirs are skipped
    try:
        import ctypes
        if not ctypes.windll.shell32.IsUserAnAdmin():
            log_fn("[WARNING] Not running as Administrator. Protected system "
                   "directories (Windows\\Prefetch, System32\\config, etc.) "
                   "will be inaccessible. Run as Admin for full coverage.", "warning")
    except Exception:
        pass  # Non-Windows or ctypes unavailable

    try:
        def _walk_onerror(e):
            if not log_fn:
                return
            msg = str(e)
            # Access Denied is expected on non-elevated runs — dim, not error
            if "Access is denied" in msg or "WinError 5" in msg:
                log_fn(f"[ACCESS DENIED] {getattr(e, 'filename', msg)}", "dim")
            else:
                log_fn(f"[WALK ERROR] {e}", "warning")

        walker = os.walk(folder, onerror=_walk_onerror)
        for root, dirs, files in walker:
            if stop_flag and stop_flag.is_set():
                if log_fn:
                    log_fn("[SCAN] Stopped during enumeration.", "warning")
                return
            # Only skip Windows metadata dirs with no forensic value ($MFT etc.)
            # Do NOT skip system dirs — keys can be saved anywhere
            dirs[:] = [d for d in dirs
                       if not d.startswith("$") or d == "$Recycle.Bin"]
            for filename in files:
                all_files.append(os.path.join(root, filename))
    except Exception as e:
        if log_fn:
            log_fn(f"[WALK FATAL] {e}", "error")
        return

    total = len(all_files)
    if log_fn:
        log_fn(f"[SCAN] Total files enumerated: {total}", "info")
        log_fn(f"[SCAN] Size filter: {_ACTIVE_MIN_BYTES} B – {_ACTIVE_MAX_BYTES} B",
               "dim")
        ext_label = ("ALL" if _ACTIVE_EXTS is None
                     else f"{len(_ACTIVE_EXTS)} selected")
        log_fn(f"[SCAN] Extension filter: {ext_label}", "dim")

    name_hits = set()

    for idx, filepath in enumerate(all_files):
        if stop_flag and stop_flag.is_set():
            if log_fn:
                log_fn("[SCAN] Stopped by user.", "warning")
            break
        if progress_fn:
            progress_fn(idx + 1, total)

        filename = os.path.basename(filepath)
        ext = os.path.splitext(filename)[1].lower()

        # ── Name-pattern search — respects extension filter AND size filter
        if do_name_search:
            if _ACTIVE_EXTS is None or ext in _ACTIVE_EXTS:
                if _size_allowed(filepath):
                    low = filename.lower()
                    for pat in BITLOCKER_NAME_PATTERNS:
                        if fnmatch.fnmatch(low, pat.lower()):
                            name_hits.add(filepath)
                            if log_fn:
                                log_fn(f"[NAME HIT] {filepath}", "success")
                            break

        # ── Content / carving path is gated by filters
        if not do_content_search and not do_binary_carve:
            continue

        # Extension filter (GUI-driven)
        if _ACTIVE_EXTS is not None and ext not in _ACTIVE_EXTS:
            continue

        # Size filter
        if not _size_allowed(filepath):
            continue

        parser = EXTENSION_PARSER_MAP.get(ext)
        if parser and do_content_search:
            try:
                parser(filepath, results, log_fn)
            except Exception as e:
                if log_fn:
                    log_fn(f"[ERROR] Parser crashed on {filepath}: {e}", "error")
        elif do_binary_carve:
            try:
                parse_raw_binary(filepath, results, log_fn)
            except Exception as e:
                if log_fn:
                    log_fn(f"[ERROR] Carve failed on {filepath}: {e}", "error")

    # Register name-only hits that produced no content match
    for path in name_hits:
        if not any(r["File Path"] == path for r in results):
            results.append({
                "File Path":       path,
                "Recovery Key ID": "Check file manually",
                "BitLocker Key":   "Suspicious filename — open to verify",
                "Validation":      "N/A",
                "Source":          "File Name Match",
            })

    # Final summary
    if log_fn:
        valid   = sum(1 for r in results if r.get("Validation") == "Valid (mod-11)")
        pattern = sum(1 for r in results if r.get("Validation") == "Pattern-only")
        log_fn(f"[SCAN COMPLETE] Total hits: {len(results)}  "
               f"|  Valid: {valid}  |  Pattern-only: {pattern}", "success")


# ═════════════════════════════════════════════════════════════════════════════
#  CSV REPORT
# ═════════════════════════════════════════════════════════════════════════════
def save_report(results, output_folder):
    """Write CSV report with fresh timestamp every invocation."""
    os.makedirs(output_folder, exist_ok=True)
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    csv_path = os.path.join(output_folder, f"BitLocker_Report_{ts}.csv")
    fieldnames = ["File Path", "Recovery Key ID", "BitLocker Key",
                  "Validation", "Source"]
    with open(csv_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        for row in results:
            # Ensure all fields present (backward compat)
            writer.writerow({k: row.get(k, "") for k in fieldnames})
    return csv_path
